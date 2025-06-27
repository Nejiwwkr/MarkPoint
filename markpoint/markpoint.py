import re

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.slide import Slide
from pptx.util import Inches, Pt

from markpoint.formula import merge_dollars_content, add_formula
from markpoint.layout import estimate_textbox_height, parse_escape_characters, complex_split


def render(file_path, target_path, target_format='pptx'):
    with open(file_path, 'r', encoding='utf-8') as f:
        content = f.read()
    if target_format == 'pptx':
        _render_ppt(content, target_path)
    pass


def _render_ppt(file_content, target_path):
    head, cover, content, back_cover = _parse_mp(file_content)
    meta_data = parse_head(head)

    #处理元数据标记的长宽比
    prs = Presentation()
    prs.slide_width = Inches(meta_data.w)
    prs.slide_height = Inches(meta_data.h)
    blank_layout = prs.slide_layouts[6]

    #封面
    slide1 = prs.slides.add_slide(blank_layout)
    parse_cover(cover, meta_data, slide1)

    #内容
    for con in content:
        parse_content(con, meta_data, prs)

    #封底
    slide2 = prs.slides.add_slide(blank_layout)
    parse_cover(back_cover, meta_data, slide2)

    prs.save(target_path)
    pass


def _parse_mp(s: str) -> tuple[str, str, list, str]:
    # 使用正则确保匹配独立的---分隔行（忽略行内---）
    parts = re.split(r'\n---\s*\n', s.strip())
    if len(parts) < 3:  # 最少需要 [配置, 封面, 封底]
        raise ValueError("Invalid Markpoint format: missing sections")

    head = parts[0]
    cover = parts[1]
    back_cover = parts[-1]
    content = parts[2:-1] if len(parts) > 3 else []
    return head, cover, content, back_cover


def _parse_hex_color(color_str: str) -> RGBColor:
    if color_str.startswith("0x"):
        hex_val = color_str[2:]
    elif color_str.startswith("#"):
        hex_val = color_str[1:]
    else:
        hex_val = color_str

    hex_val = hex_val.zfill(6)  # 处理缩写如0xfff -> 0xffffff
    return RGBColor(*bytes.fromhex(hex_val))


class MetaData:
    def __init__(self,
                 w=16,
                 h=9,
                 theme=RGBColor(0x8f, 0xaa, 0xdc),
                 background=RGBColor(222, 235, 247),
                 font="微软雅黑",
                 generate_toc=False):  # 新增目录生成选项
        self.w = w
        self.h = h
        self.theme = theme
        self.background = background
        self.font = font
        self.generate_toc = generate_toc


def parse_head(head: str) -> MetaData:
    m = MetaData()
    for line in head.splitlines():
        if '//' in line:  # 跳过注释
            line = line.split('//')[0]
        if '=' not in line:
            continue

        key, value = [part.strip() for part in line.split('=', 1)]
        value = value.split('//')[0].strip()  # 移除行尾注释

        try:
            if key in ['width', 'w']:
                m.w = float(value)
            elif key in ['height', 'h']:
                m.h = float(value)
            elif key in ['theme', 't']:
                m.theme = _parse_hex_color(value)
            elif key in ['background', 'b']:
                m.background = _parse_hex_color(value)
            elif key in ['font', 'f']:
                m.font = value
            elif key == 'generate_toc':  # 目录生成标志
                m.generate_toc = value.lower() in ['true', '1', 'yes']
        except Exception as e:
            print(f"Warning: Failed to parse config '{line}': {str(e)}")
    return m


def parse_cover(cover: str, meta_data: MetaData, slide: Slide):
    text = ''
    info = ''
    for part in cover.splitlines():
        if re.match(r'#[^#].*',part):
            text = re.sub(r'#\s*', '', part)
        if re.match(r'##[^#].*',part):
            info = re.sub(r'##\s*', '', part)

    # 添加顶层矩形
    add_rectangle(slide, meta_data.theme, [0,0,meta_data.w*0.618,meta_data.h])

    add_simple_text_box(slide, text, meta_data.font, 72, [1, meta_data.h / 2, meta_data.w - 2, 2], True)
    add_simple_text_box(slide, info, meta_data.font, 18, [meta_data.w * 0.618 + 1, meta_data.h / 2 + 2, meta_data.w * 0.322 - 1, 1])


def parse_content(content: str, meta_data: MetaData, prs: Presentation):
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)

    content = re.sub(r'//.*','',content)
    content = merge_dollars_content(content)
    current_height = 0
    spacing = 0.2
    tittle = ''

    def check_page(current, delta) -> (float, Slide):
        if current + delta > meta_data.h - 0.3:
            # 翻页
            _slide = prs.slides.add_slide(blank_layout)
            if tittle != '':
                add_rectangle(_slide, meta_data.theme, [0, 0, meta_data.w, 1.5])
                add_simple_text_box(_slide, tittle, meta_data.font, 48, [1, 0.26, meta_data.w - 2, 1.24])
                current = 1.5 + spacing
            return current, _slide
        return None

    for part in content.splitlines():
        part = part.strip()

        if re.match(r'!!!',part):
            #翻页
            slide = prs.slides.add_slide(blank_layout)
            if tittle != '':
                add_rectangle(slide, meta_data.theme, [0,0,meta_data.w,1.5])
                add_simple_text_box(slide, tittle, meta_data.font, 48, [1, 0.26, meta_data.w - 2, 1.24])
                current_height = 1.5 + spacing
        elif re.match(r'#[^#].*',part):
            tittle = re.sub(r'#\s*', '', part)
            # 页面一级标题处理
            if tittle != '':
                add_rectangle(slide, meta_data.theme, [0,0,meta_data.w,1.5])
                add_simple_text_box(slide, tittle, meta_data.font, 48, [1, 0.26, meta_data.w - 2, 1.24])
                current_height += 1.5 + spacing
        elif re.match(r'##[^#].*',part):
            # 正文一级标题
            tittleI = re.sub(r'##\s*', '', part)
            height = estimate_textbox_height(meta_data.w - 2,32,tittleI)

            if (temp := check_page(current_height, height)) is not None:
                current_height, slide = temp

            add_simple_text_box(slide, tittleI, meta_data.font, 32, [1, current_height, meta_data.w - 2, height])
            current_height += height + spacing
        elif re.match(r'###[^#].*',part):
            # 正文二级标题
            tittleII = re.sub(r'###\s*', '', part)
            height = estimate_textbox_height(meta_data.w - 2,24,tittleII)

            if (temp := check_page(current_height, height)) is not None:
                current_height, slide = temp

            add_simple_text_box(slide, tittleII, meta_data.font, 24, [1, current_height, meta_data.w - 2, height])
            current_height += height + spacing
        elif re.match(r'^\$\$.*\$\$$',part):
            # 公式
            formula = re.sub(r'\$\$', '', part)

            height, how = add_formula(slide, formula, meta_data, current_height - spacing)
            if (temp := check_page(current_height, height)) is not None:
                current_height, slide = temp
            how()

            current_height += height
        elif re.match(r'[^#\n].*',part):
            height = estimate_textbox_height(meta_data.w - 2,18,part)

            if (temp := check_page(current_height, height)) is not None:
                current_height, slide = temp

            add_simple_text_box(slide, part, meta_data.font, 18, [1, current_height, meta_data.w - 2, height])
            current_height += height + spacing

    

def add_simple_text_box(slide, text, _font, size, param, is_bold=False, color=RGBColor(0x33, 0x33, 0x33)):
    """
    easy way to add a text box

    :param slide: the target slides
    :param text: the text to be added
    :param _font: the font of the text
    :param size: the size of the text
    :param is_bold: is the text bold or not
    :param param: [left, top, width, height]
    :param color: the color of the text
    :return: nothing
    """
    #添加文本框
    textbox = slide.shapes.add_textbox(Inches(param[0]), Inches(param[1]), Inches(param[2]), Inches(param[3]))
    text_frame = textbox.text_frame
    text_frame.word_wrap = True

    # 重用第一个默认段落
    p = text_frame.paragraphs[0]

    text = text.strip()
    if "_"  not in text and "*" not in text:
        # 设置字体属性
        set_font(p, _font, size, is_bold, color)
        p.text = parse_escape_characters(text)
    else:
        text = re.sub(r'\*', '_', text)
        for e in complex_split(text, '_'):
            if e.startswith('_') and not e.startswith('__'):
                r = p.add_run()
                r.text = parse_escape_characters(e[1:])
                r.font.italic = True
            elif e.startswith('__') and not e.startswith('___'):
                r = p.add_run()
                r.text = parse_escape_characters(e[2:])
                r.font.bold = True
            elif e.startswith('___'):
                r = p.add_run()
                r.text = parse_escape_characters(e[3:])
                r.font.italic = True
                r.font.bold = True
            else:
                r = p.add_run()
                r.text = parse_escape_characters(e)
            r.font.name = _font
            r.font.size = Pt(size)
            r.font.color.rgb = color

def add_rectangle(slide, color, param):
    """
    add a rectangle

    :param slide: the target slides
    :param color: the color of the rectangle
    :param param: [left, top, width, height]
    :return:
    """
    top_rect = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(param[0]), Inches(param[1]), Inches(param[2]), Inches(param[3])
    )
    top_rect.fill.solid()
    top_rect.fill.fore_color.rgb = color


def set_font(p, _font, size, is_bold, color):
    font = p.font
    font.name = _font
    font.size = Pt(size)
    font.bold = is_bold
    font.color.rgb = color